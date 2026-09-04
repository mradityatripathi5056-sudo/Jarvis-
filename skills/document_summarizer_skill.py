"""
skills/document_summarizer_skill.py
------------------------------------------------------------
PDF, Word (.docx) aur PowerPoint (.pptx) files se text nikaal kar
existing OPENROUTER_API_KEY wale LLM se quick summary banata hai.

Extra packages chahiye (requirements-optional.txt mein add kiye hain):
    pip install pypdf python-docx python-pptx

Koi package missing ho to us file-type ke liye clear error milega,
baaki Jarvis crash nahi hoga.
"""

import json
import os

import requests
import config

MAX_CHARS_TO_LLM = 12000  # bahut lambi file ho to itna hi bhejenge (token limit)


def _extract_pdf(path: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "__MISSING__:pypdf"
    reader = PdfReader(path)
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)


def _extract_docx(path: str) -> str:
    try:
        import docx
    except ImportError:
        return "__MISSING__:python-docx"
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "__MISSING__:python-pptx"
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)


def _summarize_with_llm(text: str, instruction: str) -> str:
    try:
        resp = requests.post(
            config.OPENROUTER_URL,
            headers={
                "Authorization": f"Bearer {config.OPENROUTER_API_KEY}",
                "Content-Type": "application/json",
            },
            data=json.dumps({
                "model": config.OPENROUTER_MODEL,
                "messages": [
                    {
                        "role": "user",
                        "content": f"{instruction}\n\nDocument content:\n{text[:MAX_CHARS_TO_LLM]}",
                    }
                ],
                "max_tokens": 700,
            }),
            timeout=40,
        )
        data = resp.json()
        return data["choices"][0]["message"]["content"].strip()
    except Exception as e:
        return f"Summary banane mein LLM error aaya: {e}"


def summarize_document(params: dict) -> str:
    path = params.get("path", "").strip()
    instruction = params.get(
        "instruction",
        "Is document ka short, clear summary Hinglish mein do - "
        "main points bullet mein, aur agar koi action item ho to wo bhi.",
    )
    if not path or not os.path.exists(path):
        return f"'{path}' file nahi mili. Poora path batao."

    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        text = _extract_pdf(path)
    elif ext == ".docx":
        text = _extract_docx(path)
    elif ext == ".pptx":
        text = _extract_pptx(path)
    else:
        return f"'{ext}' file type supported nahi hai - sirf .pdf, .docx, .pptx chalti hain."

    if text.startswith("__MISSING__:"):
        package = text.split(":", 1)[1]
        return f"Is file type ke liye package missing hai. Chalao: pip install {package}"

    if not text.strip():
        return "File se koi text nikal nahi paya (shayad scanned/image-based document hai - OCR chahiye hoga)."

    return _summarize_with_llm(text, instruction)


ACTIONS = {
    "summarize_document": summarize_document,
}

DOCS = """
- summarize_document: {"path": "report.pdf"}
    (PDF/.docx/.pptx file ka text nikal kar LLM se short summary banata hai;
    "instruction" param optional hai agar summary ka style/focus batana ho)

Example:
User: "is PDF ka summary bata do - C:\\Users\\me\\report.pdf"
-> {"actions": [{"action": "summarize_document", "params": {"path": "C:\\Users\\me\\report.pdf"}}]}

User: "ye PPT mein sirf financial points nikaal ke batao"
-> {"actions": [{"action": "summarize_document", "params": {"path": "<file path>", "instruction": "Sirf financial/numbers related points nikaal ke batao"}}]}
"""
